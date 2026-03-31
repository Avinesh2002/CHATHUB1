import os
import io
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text(file_stream, filename):
    """
    Extracts text from a given file stream based on its extension.
    Supported extensions: .txt, .pdf, .pptx, .ppt
    """
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    
    try:
        if ext == ".txt":
            # For txt, decode the byte stream
            text = file_stream.read().decode('utf-8', errors='ignore')
            
        elif ext == ".pdf":
            # For PDF, use PyPDF2
            pdf_reader = PdfReader(io.BytesIO(file_stream.read()))
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                    
        elif ext in [".ppt", ".pptx"]:
            # For PPT/PPTX
            prs = Presentation(io.BytesIO(file_stream.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            raise ValueError(f"Unsupported file format: {ext}")
            
    except Exception as e:
        print(f"Error processing {filename}: {str(e)}")
        raise e
        
    return text.strip()
